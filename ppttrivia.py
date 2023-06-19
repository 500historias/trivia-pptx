import random
import pandas as pd
from pptx import Presentation
from pptx import dml
from pptx.enum.text import MSO_AUTO_SIZE 

# Read the Excel file
df = pd.read_excel('./templates/questions.xlsx')  # questions must be in the first sheet of the document

# Randomize the rows
df = df.sample(frac=1).reset_index(drop=True)

# Load the existing PowerPoint template
template_path = './templates/triviappt.pptx'
presentation = Presentation(template_path)

# Define the slide layout index for question slides
question_layout_index = 0

# Iterate over each row in the DataFrame
for _, row in df.iterrows():
    # Extract question and answers from the row
    question = row[0]
    answers = random.sample(row[1:].tolist(), k=4)   # radomize order of answers

    # Create the first slide with question and random answers
    slide = presentation.slides.add_slide(presentation.slide_layouts[question_layout_index])
    slide.shapes.title.text = question

    # Add the random answers to the slide
    for i, answer in enumerate(answers):
        j =  i+10
        content = slide.placeholders[j].text_frame
        content.text = answer

    # Create the second slide with the correct answer highlighted
    slide = presentation.slides.add_slide(presentation.slide_layouts[question_layout_index])
    slide.shapes.title.text = question

    # Add the answers to the slide, highlighting the correct answer
    for i, answer in enumerate(answers):
        j = i+10
        content = slide.placeholders[j].text_frame
        content.text = answer
        content.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        if answer == row[1]:  # Assuming the correct answer is in the second column
            # content.text = answer
            # content.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            slide.placeholders[j].fill.solid()
            slide.placeholders[j].fill.fore_color.rgb = dml.color.RGBColor.from_string("EB6924")


# Save the modified PowerPoint presentation
output_path = './output/output_slides.pptx'
presentation.save(output_path)
