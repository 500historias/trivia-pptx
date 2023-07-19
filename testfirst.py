import random
import pandas as pd
from pptx import Presentation

# Load the existing PowerPoint template
template_path = './templates/triviappt99.pptx'
presentation = Presentation(template_path)

slide = presentation.slides.add_slide(presentation.slide_layouts[0])

for shape in slide.placeholders:
    print ('%d %s [%s]' % (shape.placeholder_format.idx, shape.name, shape.placeholder_format.type))